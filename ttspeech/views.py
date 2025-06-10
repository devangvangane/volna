import os
from PIL import Image
from pptx import Presentation
import pytesseract as pytesseract
import pyttsx3
import docx
import PyPDF2
from django.shortcuts import render
from django.http import JsonResponse
from django.conf import settings
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage
import requests
from gtts import gTTS
import json
from googletrans import Translator
from gtts.lang import tts_langs

# Media folder for audio storage
MEDIA_FOLDER = os.path.join(settings.BASE_DIR, "media")
os.makedirs(MEDIA_FOLDER, exist_ok=True)
GEMINI_API_KEY = "AIzaSyAHygzhpSGaLo2ZiqkoLeORfa_a9iugEf0"
GEMINI_API_URL = "https://generativelanguage.googleapis.com/v1/models/gemini-1.5-flash:generateContent"

# Initialize TTS Engine
# engine = pyttsx3.init()

# Fetch available voices
# voices = engine.getProperty("voices")
VOICE_MAP = {
    "default": None,  # Uses system default voice
    "male": 0,        # First available voice
    "female": 1,      # Second available voice (if available)
}

def home(request):
    """Render homepage with all generated files."""
    functionality = request.GET.get("functionality", "tts")
    files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
    return render(request, "index.html", {"functionality": functionality, "files": files})

def improve_text_with_gemini_grammar(text):
    """
    Use Gemini API to correct grammar and clarity of user-entered text.
    """
    payload = {
        "contents": [
            {
                "parts": [
                    {
                        "text": f"Please correct the grammar and improve clarity of this text:\n\n{text}"
                    }
                ]
            }
        ]
    }

    headers = {"Content-Type": "application/json"}

    try:
        response = requests.post(
            f"{GEMINI_API_URL}?key={GEMINI_API_KEY}",
            headers=headers,
            data=json.dumps(payload)
        )

        if response.status_code == 200:
            response_data = response.json()
            if 'candidates' in response_data and len(response_data['candidates']) > 0:
                content = response_data['candidates'][0].get('content', {})
                if 'parts' in content and len(content['parts']) > 0:
                    return content['parts'][0].get('text', text)
        print(f"Gemini Grammar API error: {response.status_code} - {response.text}")
        return text
    except Exception as e:
        print(f"Gemini grammar correction error: {str(e)}")
        return text


@csrf_exempt
def text_to_speech(request):
    """Convert text input to speech with optional grammar correction using Gemini AI."""
    engine = pyttsx3.init()
    voices = engine.getProperty("voices")
    if request.method == "POST":
        text = request.POST.get("text", "").strip()
        speed = float(request.POST.get("speed", 150))
        voice = request.POST.get("voice", "default")
        correct_grammar = request.POST.get("correct_grammar", "false").lower() == "true"

        if not text:
            return JsonResponse({"error": "Text field is required."}, status=400)

        if correct_grammar:
            print("Applying grammar correction via Gemini")
            text = improve_text_with_gemini_grammar(text)
            print("Corrected text:", text)

        # Set voice properties
        engine.setProperty("rate", speed)

        if voice in VOICE_MAP and VOICE_MAP[voice] is not None and VOICE_MAP[voice] < len(voices):
            engine.setProperty("voice", voices[VOICE_MAP[voice]].id)

        file_name = f"speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
        file_path = os.path.join(MEDIA_FOLDER, file_name)

        engine.save_to_file(text, file_path)
        engine.runAndWait()

        files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
        return JsonResponse({"audio_url": f"/media/{file_name}", "files": files})
@csrf_exempt
def file_to_speech(request):
    engine = pyttsx3.init()
    voices = engine.getProperty("voices")
    """Convert uploaded PDF/DOCX to speech and return updated file list."""
    if request.method == "POST" and request.FILES.get("file"):
        uploaded_file = request.FILES["file"]
        file_extension = uploaded_file.name.split(".")[-1].lower()
        speed = float(request.POST.get("speed", 150))
       
        voice = request.POST.get("voice", "default")

        file_path = os.path.join(MEDIA_FOLDER, uploaded_file.name)
        with default_storage.open(file_path, "wb+") as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)

        extracted_text = ""

        try:
            if file_extension == "pdf":
                with open(file_path, "rb") as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        extracted_text += page.extract_text() + " "
            elif file_extension == "docx":
                doc = docx.Document(file_path)
                extracted_text = " ".join([para.text for para in doc.paragraphs])
            else:
                return JsonResponse({"error": "Unsupported file type."}, status=400)

            if not extracted_text.strip():
                return JsonResponse({"error": "No text found in file."}, status=400)

            # Set TTS properties
            engine.setProperty("rate", speed)
         
            print(VOICE_MAP)
            if voice in VOICE_MAP and VOICE_MAP[voice] is not None and VOICE_MAP[voice] < len(voices):
                engine.setProperty("voice", voices[VOICE_MAP[voice]].id)

            audio_filename = f"file_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            audio_filepath = os.path.join(MEDIA_FOLDER, audio_filename)

            engine.save_to_file(extracted_text, audio_filepath)
            engine.runAndWait()

            files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
            return JsonResponse({"audio_url": f"/media/{audio_filename}", "files": files})

        except Exception as e:
            return JsonResponse({"error": f"Error: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)


@csrf_exempt
def image_to_speech(request):
    engine = pyttsx3.init()
    voices = engine.getProperty("voices")
    """Extract text from an image and convert it to speech."""
    if request.method == "POST" and request.FILES.get("image"):
        image_file = request.FILES["image"]

        # Save uploaded image
        image_path = os.path.join(MEDIA_FOLDER, image_file.name)
        with default_storage.open(image_path, "wb+") as destination:
            for chunk in image_file.chunks():
                destination.write(chunk)

        try:
            # Open and process the image
            image = Image.open(image_path)
            extracted_text = pytesseract.image_to_string(image)

            if not extracted_text.strip():
                return JsonResponse({"error": "No readable text found in the image."}, status=400)

            # Generate speech file
            audio_filename = f"image_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            audio_filepath = os.path.join(MEDIA_FOLDER, audio_filename)

            engine.save_to_file(extracted_text, audio_filepath)
            engine.runAndWait()
            print(audio_filepath)
            return JsonResponse({"audio_url": f"/media/{audio_filename}"})

        except Exception as e:
            return JsonResponse({"error": f"Error: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)


# Gemini API Config


def improve_text_with_gemini(text):
    """Send extracted text to Gemini API and get an improved version."""
    payload = {
        "contents": [
            {
                "parts": [
                    {
                        "text": f"this is a ppt words that extract from ppt convert into speech:\n\n{text}"
                    }
                ]
            }
        ]
    }

    headers = {"Content-Type": "application/json"}

    try:
        response = requests.post(
            f"{GEMINI_API_URL}?key={GEMINI_API_KEY}",
            headers=headers,
            data=json.dumps(payload)
        )

        if response.status_code == 200:
            response_data = response.json()
            # Properly extract text from the Gemini API response structure
            if 'candidates' in response_data and len(response_data['candidates']) > 0:
                if 'content' in response_data['candidates'][0]:
                    content = response_data['candidates'][0]['content']
                    if 'parts' in content and len(content['parts']) > 0:
                        return content['parts'][0].get('text', text)
            # Log the response for debugging
            print(f"Gemini API response: {response_data}")
        else:
            print(f"Gemini API error: {response.status_code} - {response.text}")

        return text  # Fallback to original if parsing fails
    except Exception as e:
        print(f"Error calling Gemini API: {str(e)}")
        return text  # Fallback to original if API call fails


@csrf_exempt
def pptx_to_speech(request):
    engine = pyttsx3.init()
    voices = engine.getProperty("voices")
    """Extract text from a PPTX file, improve it with Gemini, and convert it to speech."""
    if request.method == "POST" and request.FILES.get("ppt"):
        pptx_file = request.FILES["ppt"]
        speed = float(request.POST.get("speed", 150))
  
        voice = request.POST.get("voice", "default")

        if not pptx_file.name.lower().endswith(".pptx"):
            return JsonResponse({"error": "Only .pptx files are supported."}, status=400)

        # Save the uploaded PPTX file
        file_path = os.path.join(MEDIA_FOLDER, pptx_file.name)
        with default_storage.open(file_path, "wb+") as destination:
            for chunk in pptx_file.chunks():
                destination.write(chunk)

        try:
            # Parse presentation and extract text
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        for paragraph in shape.text_frame.paragraphs:
                            full_line = ''.join([run.text for run in paragraph.runs]).strip()
                            if full_line:
                                text_runs.append(full_line)

            extracted_text = "\n".join(text_runs)

            if not extracted_text.strip():
                return JsonResponse({"error": "No readable text found in the presentation."}, status=400)

            # Improve text with Gemini API
            improved_text = improve_text_with_gemini(extracted_text)
            print(improved_text)
            # Add debug logging
            print(f"Original text length: {len(extracted_text)}")
            print(f"Improved text length: {len(improved_text)}")

            # Set TTS engine properties
            engine.setProperty("rate", speed)
           

            # Handle voice selection more safely
            voice_index = VOICE_MAP.get(voice)
            if voice_index is not None and isinstance(voice_index, int) and 0 <= voice_index < len(voices):
                engine.setProperty("voice", voices[voice_index].id)
            else:
                # Default to first voice if specified voice is invalid
                if len(voices) > 0:
                    engine.setProperty("voice", voices[0].id)

            audio_filename = f"pptx_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            audio_filepath = os.path.join(MEDIA_FOLDER, audio_filename)

            # Generate speech from improved text
            engine.save_to_file(improved_text, audio_filepath)
            engine.runAndWait()

            # Make sure the file was created
            if not os.path.exists(audio_filepath):
                return JsonResponse({"error": "Failed to generate audio file."}, status=500)

            print(f"Audio file created: {audio_filepath}")

            files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
            return JsonResponse({"audio_url": f"/media/{audio_filename}", "files": files})

        except Exception as e:
            import traceback
            print(f"Error processing PPTX: {str(e)}")
            print(traceback.format_exc())
            return JsonResponse({"error": f"Error processing PPTX: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)

SUPPORTED_LANGUAGES = [
    "af", "ar", "bg", "bn", "bs", "ca", "cs", "cy", "da", "de", "el", "en", 
    "eo", "es", "et", "fi", "fr", "gu", "hi", "hr", "hu", "hy", "id", "is", 
    "it", "iw", "ja", "jw", "km", "kn", "ko", "la", "lv", "mk", "ml", "mr", 
    "ms", "my", "ne", "nl", "no", "pl", "pt", "ro", "ru", "si", "sk", "sq", 
    "sr", "su", "sv", "sw", "ta", "te", "th", "tl", "tr", "uk", "ur", "vi", 
    "zh-cn", "zh-tw", "zu"
]

def get_supported_languages(request):
    """Returns a list of supported languages for gTTS."""
    return JsonResponse({"languages": tts_langs()})


@csrf_exempt
def translation(request):
    """Translate text and convert it to speech."""
    if request.method == "POST":
        text = request.POST.get("text")
        target_lang = request.POST.get("language", "en")  # Default is English

        if not text.strip():
            return JsonResponse({"error": "Text field is required."}, status=400)

        if target_lang not in SUPPORTED_LANGUAGES:
            return JsonResponse({"error": f"Unsupported language: {target_lang}"}, status=400)

        try:
            # Step 1: Detect source language (if not specified)
            translator = Translator()
            detected = translator.detect(text)
            source_lang = detected.lang
            
            # Step 2: Translate the text
            translated_text = translator.translate(text, src=source_lang, dest=target_lang).text

            # Step 3: Convert translated text to speech
            file_name = f"translated_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            file_path = os.path.join(MEDIA_FOLDER, file_name)

            tts = gTTS(text=translated_text, lang=target_lang, slow=False)
            tts.save(file_path)

            return JsonResponse({
                "source_language": source_lang,
                "translated_text": translated_text,
                "audio_url": f"/media/{file_name}",
                "files": os.listdir(MEDIA_FOLDER)  # Send updated file list
            })

        except Exception as e:
            return JsonResponse({"error": f"Error: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)
@csrf_exempt
def delete_audio(request):
    """Delete selected audio file and update the file list."""
    if request.method == "POST":
        file_name = request.POST.get("file_name", "").strip()

        if not file_name:
            return JsonResponse({"error": "File name is required."}, status=400)

        file_path = os.path.join(MEDIA_FOLDER, file_name)

        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
                return JsonResponse({"message": "File deleted successfully.", "files": files})
            except Exception as e:
                return JsonResponse({"error": f"Error deleting file: {str(e)}"}, status=500)
        else:
            return JsonResponse({"error": "File not found."}, status=404)

    return JsonResponse({"error": "Invalid request."}, status=400)
